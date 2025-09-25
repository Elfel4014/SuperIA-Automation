from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Pt

class PowerPointAutomator:
    def __init__(self):
        pass

    def create_presentation(self, filename="apresentacao_superia.pptx", title="Apresentação SuperIA", subtitle="Automação de PC com Controle Total"):
        try:
            prs = Presentation()
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title_placeholder = slide.shapes.title
            subtitle_placeholder = slide.placeholders[1]

            title_placeholder.text = title
            subtitle_placeholder.text = subtitle

            prs.save(filename)
            print(f"Apresentação PowerPoint \'{filename}\' criada com sucesso.")
            return True
        except Exception as e:
            print(f"Erro ao criar apresentação PowerPoint: {e}")
            return False

    def add_title_and_content_slide(self, filename, slide_title, content_points):
        try:
            prs = Presentation(filename)
            bullet_slide_layout = prs.slide_layouts[1] # Layout de Título e Conteúdo
            slide = prs.slides.add_slide(bullet_slide_layout)

            title_placeholder = slide.shapes.title
            body_placeholder = slide.placeholders[1]

            title_placeholder.text = slide_title

            tf = body_placeholder.text_frame
            tf.clear() # Limpa qualquer texto existente
            for point in content_points:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0 # Nível do bullet point

            prs.save(filename)
            print(f"Slide de título e conteúdo adicionado à apresentação \'{filename}\' com sucesso.")
            return True
        except Exception as e:
            print(f"Erro ao adicionar slide de título e conteúdo: {e}")
            return False

    def add_picture_slide(self, filename, slide_title, image_path, left=Inches(1), top=Inches(2), width=Inches(8)):
        try:
            prs = Presentation(filename)
            blank_slide_layout = prs.slide_layouts[6] # Layout em branco
            slide = prs.slides.add_slide(blank_slide_layout)

            # Adicionar título ao slide
            left_title = Inches(0.5)
            top_title = Inches(0.5)
            width_title = Inches(9)
            height_title = Inches(1)
            title_shape = slide.shapes.add_textbox(left_title, top_title, width_title, height_title)
            title_frame = title_shape.text_frame
            title_frame.text = slide_title
            title_frame.paragraphs[0].font.size = Pt(36)
            title_frame.paragraphs[0].font.bold = True

            # Adicionar imagem
            slide.shapes.add_picture(image_path, left, top, width=width)

            prs.save(filename)
            print(f"Slide de imagem adicionado à apresentação \'{filename}\' com sucesso.")
            return True
        except Exception as e:
            print(f"Erro ao adicionar slide de imagem: {e}")
            return False

if __name__ == "__main__":
    automator = PowerPointAutomator()
    test_pptx = "test_superia_presentation.pptx"

    print("\n--- Criando apresentação PowerPoint ---")
    automator.create_presentation(test_pptx, "Demonstração SuperIA", "Capacidades de Automação")

    print("\n--- Adicionando slide de título e conteúdo ---")
    content = [
        "Manipulação do Registro do Windows",
        "Shell Scripting Linux Avançado",
        "Automação Office Completa",
        "Gerenciamento de Processos do Sistema",
        "Integração de APIs do Sistema"
    ]
    automator.add_title_and_content_slide(test_pptx, "Funcionalidades Principais", content)

    # Para testar add_picture_slide, você precisaria de um arquivo de imagem existente
    # Exemplo: criar um arquivo de imagem dummy
    try:
        from PIL import Image
        img = Image.new('RGB', (60, 30), color = 'red')
        img.save('dummy_image.png')
        print("Imagem dummy criada: dummy_image.png")
        automator.add_picture_slide(test_pptx, "Exemplo de Imagem", "dummy_image.png")
        os.remove('dummy_image.png')
        print("Imagem dummy removida.")
    except ImportError:
        print("Pillow não está instalada. Não foi possível criar imagem dummy para teste.")
        print("Para testar o slide de imagem, instale Pillow (pip install Pillow) e forneça um caminho de imagem válido.")
    except Exception as e:
        print(f"Erro ao criar ou adicionar imagem dummy: {e}")

    print(f"\nVerifique o arquivo \'{test_pptx}\' para ver o resultado.")

